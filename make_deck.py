"""
Tendly pitch deck — premium rebuild
Run: python3 make_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ────────────────────────────────────────────────────────────────────
INK    = RGBColor(0x09, 0x09, 0x0F)   # near-black with a whisper of indigo
INK2   = RGBColor(0x11, 0x11, 0x1E)   # slightly lighter for cards
GOLD   = RGBColor(0xC8, 0xA9, 0x6E)   # champagne gold
GOLD2  = RGBColor(0xE8, 0xD4, 0xA8)   # pale gold for body text
CREAM  = RGBColor(0xF5, 0xF0, 0xEB)   # off-white
MIST   = RGBColor(0x55, 0x55, 0x6E)   # muted blue-grey for captions
SAGE   = RGBColor(0x6B, 0x9B, 0x8E)   # muted teal-sage accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DIM    = RGBColor(0x33, 0x33, 0x44)   # dim for dividers/rules

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ── XML helpers ────────────────────────────────────────────────────────────────

def rgb(c: RGBColor) -> str:
    return str(c)

def _spPr(shape):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    return spPr

def no_line(shape):
    spPr = _spPr(shape)
    ln = spPr.find(qn('a:ln'))
    if ln is not None: spPr.remove(ln)
    ln = etree.SubElement(spPr, qn('a:ln'))
    etree.SubElement(ln, qn('a:noFill'))

def solid_fill(shape, color: RGBColor):
    spPr = _spPr(shape)
    for tag in ['a:noFill','a:solidFill','a:gradFill','a:pattFill','a:blipFill']:
        el = spPr.find(qn(tag))
        if el is not None: spPr.remove(el)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    clr = etree.SubElement(sf, qn('a:srgbClr'))
    clr.set('val', rgb(color))

def grad_fill(shape, c1: RGBColor, c2: RGBColor, ang=5400000):
    spPr = _spPr(shape)
    for tag in ['a:noFill','a:solidFill','a:gradFill','a:pattFill','a:blipFill']:
        el = spPr.find(qn(tag))
        if el is not None: spPr.remove(el)
    xml = f'''<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:gsLst>
    <a:gs pos="0"><a:srgbClr val="{rgb(c1)}"/></a:gs>
    <a:gs pos="100000"><a:srgbClr val="{rgb(c2)}"/></a:gs>
  </a:gsLst>
  <a:lin ang="{ang}" scaled="0"/>
</a:gradFill>'''
    spPr.append(etree.fromstring(xml))

def round_corners(shape, val=60000):
    spPr = _spPr(shape)
    pg = spPr.find(qn('a:prstGeom'))
    if pg is None:
        pg = etree.SubElement(spPr, qn('a:prstGeom'))
    pg.set('prst', 'roundRect')
    av = pg.find(qn('a:avLst'))
    if av is None: av = etree.SubElement(pg, qn('a:avLst'))
    gd = av.find(qn('a:gd'))
    if gd is None: gd = etree.SubElement(av, qn('a:gd'))
    gd.set('name','adj'); gd.set('fmla', f'val {val}')

def shadow(shape):
    spPr = _spPr(shape)
    xml = '''<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <a:outerShdw blurRad="60000" dist="30000" dir="5400000" rotWithShape="0">
    <a:srgbClr val="000000"><a:alpha val="30000"/></a:srgbClr>
  </a:outerShdw>
</a:effectLst>'''
    spPr.append(etree.fromstring(xml))

def alpha(shape, pct):
    spPr = _spPr(shape)
    val = str(int(pct * 1000))
    for fill_tag in [qn('a:solidFill'), qn('a:gradFill')]:
        fill = spPr.find(fill_tag)
        if fill is not None:
            for clr in fill.iter(qn('a:srgbClr')):
                a_el = etree.SubElement(clr, qn('a:alpha'))
                a_el.set('val', val)

def box(sl, x, y, w, h, color=None, r=None):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    no_line(s)
    if color: solid_fill(s, color)
    if r: round_corners(s, r)
    return s

def gbox(sl, x, y, w, h, c1, c2, ang=5400000, r=None):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    no_line(s); grad_fill(s, c1, c2, ang)
    if r: round_corners(s, r)
    return s

def circ(sl, cx, cy, r, color):
    s = sl.shapes.add_shape(9, Inches(cx-r), Inches(cy-r), Inches(r*2), Inches(r*2))
    no_line(s); solid_fill(s, color)
    return s

def txt(sl, text, x, y, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Helvetica Neue'
    return tb

def txt2(sl, lines, x, y, w, h, align=PP_ALIGN.LEFT):
    """lines = [(text, size, bold, italic, color)]"""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for (t, sz, bd, it, col) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = t
        run.font.size = Pt(sz)
        run.font.bold = bd
        run.font.italic = it
        run.font.color.rgb = col
        run.font.name = 'Helvetica Neue'
    return tb

def rule(sl, x, y, w, color=GOLD, h=0.025):
    b = box(sl, x, y, w, h, color)
    return b

# ══════════════════════════════════════════════════════════════════════════════
#  01 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)

box(sl, 0, 0, 13.33, 7.5, INK)

# subtle radial glow top-right
g1 = circ(sl, 12.5, 0.8, 4.0, GOLD); alpha(g1, 6)
g2 = circ(sl, 12.8, 0.5, 2.0, GOLD); alpha(g2, 8)

# thin gold rule left edge
rule(sl, 0.9, 2.15, 1.8)

txt(sl, "tendly", 0.9, 2.45, 9, 1.8, size=96, bold=True, color=WHITE)
txt(sl, "AI-Powered Care for Every Voice",
    0.9, 4.35, 8, 0.65, size=22, color=GOLD2)

# bottom label
txt(sl, "Hackathon Demo  ·  2025",
    0.9, 6.85, 6, 0.38, size=11, color=MIST)

# ══════════════════════════════════════════════════════════════════════════════
#  02 — PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
box(sl, 0, 0, 13.33, 7.5, INK)

txt(sl, "Problem", 0.9, 0.55, 4, 0.45, size=11, color=GOLD)
rule(sl, 0.9, 1.1, 0.6)
txt(sl, "Seniors can't get\nhelp when they\nneed it most.",
    0.9, 1.25, 7, 2.8, size=52, bold=True, color=WHITE)

# three facts — right column
facts = [
    ("54 million", "Americans aged 65+ need daily care"),
    ("28%",        "Of adverse events stem from missed communications"),
    ("10,000",     "Baby boomers turn 65 every single day"),
]
for i, (stat, label) in enumerate(facts):
    fy = 1.4 + i * 1.85
    b = box(sl, 8.5, fy, 4.35, 1.55, INK2, r=30000); shadow(b)
    txt(sl, stat,  8.7, fy+0.1,  3.9, 0.7,  size=34, bold=True, color=GOLD)
    txt(sl, label, 8.7, fy+0.82, 3.9, 0.6,  size=12.5, color=GOLD2)

rule(sl, 0, 7.4, 13.33, DIM, h=0.04)

# ══════════════════════════════════════════════════════════════════════════════
#  03 — SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
box(sl, 0, 0, 13.33, 7.5, INK)

txt(sl, "Solution", 0.9, 0.55, 4, 0.45, size=11, color=GOLD)
rule(sl, 0.9, 1.1, 0.6)
txt(sl, "One platform.\nThree superpowers.",
    0.9, 1.25, 9, 1.7, size=52, bold=True, color=WHITE)

cards = [
    ("Patient\nVoice",       "One tap. Speak any need. Instantly routed.",         SAGE),
    ("Caregiver\nDashboard", "Live feed, urgency-ranked by AI.",                   GOLD),
    ("AI Triage\n& Action",  "Claude classifies. Simular executes.",               RGBColor(0x8B,0x5C,0xF6)),
]
for i, (title, sub, accent) in enumerate(cards):
    cx = 0.9 + i * 4.1
    c = box(sl, cx, 3.3, 3.75, 3.7, INK2, r=40000); shadow(c)
    rule(sl, cx+0.3, 3.55, 0.5, accent)
    txt(sl, title, cx+0.3, 3.78, 3.2, 0.95, size=21, bold=True, color=WHITE)
    txt(sl, sub,   cx+0.3, 4.82, 3.2, 0.9,  size=13.5, color=GOLD2)

# ══════════════════════════════════════════════════════════════════════════════
#  04 — HOW IT WORKS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
box(sl, 0, 0, 13.33, 7.5, INK)

txt(sl, "How It Works", 0.9, 0.55, 6, 0.45, size=11, color=GOLD)
rule(sl, 0.9, 1.1, 0.6)
txt(sl, "Request to resolution\nin under 2 seconds.",
    0.9, 1.25, 9, 1.5, size=44, bold=True, color=WHITE)

steps = [
    ("01", "Patient Speaks",   "Voice or text"),
    ("02", "Deepgram STT",     "< 300ms"),
    ("03", "Claude Triage",    "Category + urgency"),
    ("04", "Alert",            "Live dashboard push"),
    ("05", "Simular Executes", "Task done"),
]
SW, SH = 2.25, 2.55
for i, (num, title, sub) in enumerate(steps):
    sx = 0.6 + i * 2.52
    c = box(sl, sx, 3.3, SW, SH, INK2, r=30000); shadow(c)
    txt(sl, num,   sx+0.2, 3.42, SW-0.3, 0.45, size=11, color=GOLD)
    txt(sl, title, sx+0.2, 3.85, SW-0.3, 0.65, size=15, bold=True, color=WHITE)
    txt(sl, sub,   sx+0.2, 4.58, SW-0.3, 0.55, size=12, color=MIST)
    if i < len(steps)-1:
        txt(sl, "→", sx+SW+0.1, 4.35, 0.32, 0.4, size=16, color=DIM, align=PP_ALIGN.CENTER)

# supporting layer labels
layers = ["Redis memory enriches every step",
          "Arize logs every AI decision",
          "Sentry monitors in production"]
for i, lb in enumerate(layers):
    txt(sl, "— " + lb, 0.9, 6.3 + i*0.35, 12, 0.33, size=11, color=MIST)

# ══════════════════════════════════════════════════════════════════════════════
#  05 — PATIENT INTERFACE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
box(sl, 0, 0, 13.33, 7.5, INK)

txt(sl, "Patient", 0.9, 0.55, 6, 0.45, size=11, color=GOLD)
rule(sl, 0.9, 1.1, 0.6)
txt(sl, "Simple enough for\nanyone.",
    0.9, 1.25, 6, 1.5, size=52, bold=True, color=WHITE)
txt(sl, "One button. Any request. Instant response.",
    0.9, 2.85, 6, 0.55, size=16, color=GOLD2)

features = [
    "Voice recording — no typing required",
    "Audio talkback — AI speaks the response",
    "Confirmation gate for sensitive actions",
    "Recent request history at a glance",
]
for i, f in enumerate(features):
    fy = 3.75 + i * 0.7
    rule(sl, 0.9, fy+0.15, 0.18, SAGE, h=0.32)
    txt(sl, f, 1.22, fy, 5.3, 0.6, size=14, color=CREAM)

# Right: minimal phone wireframe
ph = box(sl, 8.2, 0.45, 4.5, 6.8, INK2, r=100000); shadow(ph)
rule(sl, 8.2, 0.45, 4.5, DIM, h=0.04)

# notch
box(sl, 9.95, 0.46, 0.9, 0.22, INK, r=50000)

# screen zones
box(sl, 8.45, 0.85, 2.0, 5.85, RGBColor(0x0D,0x0D,0x18))  # left dark strip
box(sl, 10.45, 0.85, 2.0, 5.85, RGBColor(0xF0,0xED,0xE8))  # right light strip

txt(sl, "Good morning,\nMary.", 8.5, 1.1, 1.85, 0.85, size=10, bold=True, color=WHITE)
txt(sl, "Your care team", 8.5, 2.15, 1.85, 0.3, size=8, color=MIST)

for ii, (ini, nm) in enumerate([("JP","Jenny P."), ("CM","Carlos M.")]):
    ty = 2.52 + ii*0.5
    av = box(sl, 8.52, ty, 0.32, 0.32, SAGE, r=50000)
    txt(sl, ini, 8.52, ty+0.03, 0.32, 0.24, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, nm, 8.88, ty+0.05, 1.5, 0.25, size=8, color=GOLD2)

# mic button
mic = box(sl, 10.65, 2.0, 1.5, 1.5, SAGE, r=100000); shadow(mic)
txt(sl, "MIC", 10.65, 2.45, 1.5, 0.6, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "What do you need?", 10.5, 1.6, 1.85, 0.35, size=9, color=RGBColor(0x44,0x44,0x55), align=PP_ALIGN.CENTER)
txt(sl, "Tap to speak", 10.5, 3.6, 1.85, 0.35, size=8.5, color=MIST, align=PP_ALIGN.CENTER)

# activity items
txt(sl, "Recent", 10.55, 4.15, 1.8, 0.3, size=8, bold=True, color=RGBColor(0x33,0x33,0x44))
for ii, (label, st, col) in enumerate([
    ("Blanket requested","Done", GREEN := RGBColor(0x22,0xC5,0x5E)),
    ("Dizziness report", "Active", RGBColor(0xEF,0x44,0x44)),
    ("Sinatra playing",  "Done",  RGBColor(0x22,0xC5,0x5E)),
]):
    ay = 4.55 + ii * 0.62
    ac = box(sl, 10.48, ay, 1.9, 0.5, WHITE := RGBColor(0xFF,0xFF,0xFF), r=20000); shadow(ac)
    dot = box(sl, 10.52, ay+0.17, 0.12, 0.12, col, r=50000)
    txt(sl, label, 10.68, ay+0.06, 1.6, 0.25, size=7.5, color=RGBColor(0x22,0x22,0x33))
    txt(sl, st,    10.68, ay+0.29, 1.6, 0.2,  size=7,   color=MIST)

# ══════════════════════════════════════════════════════════════════════════════
#  06 — CAREGIVER DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
box(sl, 0, 0, 13.33, 7.5, INK)

txt(sl, "Caregiver", 0.9, 0.55, 6, 0.45, size=11, color=GOLD)
rule(sl, 0.9, 1.1, 0.6)
txt(sl, "Every request.\nEvery priority.\nIn real time.",
    0.9, 1.25, 5.5, 2.4, size=44, bold=True, color=WHITE)

# Stat pills
for i, (val, lbl) in enumerate([("4","Active"),("1","Emergency"),("7","Resolved")]):
    px = 0.9 + i * 1.7
    pb = box(sl, px, 3.95, 1.45, 0.9, INK2, r=30000)
    txt(sl, val, px, 3.98, 1.45, 0.45, size=22, bold=True,
        color=RGBColor(0xEF,0x44,0x44) if lbl=="Emergency" else GOLD, align=PP_ALIGN.CENTER)
    txt(sl, lbl, px, 4.42, 1.45, 0.38, size=9.5, color=MIST, align=PP_ALIGN.CENTER)

# Request cards — right side
card_data = [
    (RGBColor(0xEF,0x44,0x44), "Emergency · Medical",   "Margaret Chen  Rm 14",
     "Chest tightness, difficulty breathing.",  "→ Respond immediately"),
    (RGBColor(0xF9,0x73,0x16), "High · Medical",        "Robert Kim  Rm 22",
     "Dizziness following new medication.",     "→ Assess soon"),
    (SAGE,                     "Low · Family Comms",    "Dorothy Walsh  Rm 8",
     "Wants to call her daughter.",             "→ Initiate call (confirm)"),
    (GOLD,                     "Low · Automated",       "James Patel  Rm 31",
     "Wants to listen to Frank Sinatra.",       "→ Run automation"),
]
for i, (accent, tag, patient, summary, action) in enumerate(card_data):
    row, col = divmod(i, 2)
    cx = 6.3 + col * 3.48
    cy = 0.55 + row * 3.4
    c = box(sl, cx, cy, 3.25, 3.05, INK2, r=30000); shadow(c)
    rule(sl, cx, cy, 3.25, accent, h=0.04)
    txt(sl, tag,     cx+0.25, cy+0.22, 2.75, 0.35, size=10,   color=accent)
    txt(sl, patient, cx+0.25, cy+0.62, 2.75, 0.38, size=13.5, bold=True, color=WHITE)
    txt(sl, summary, cx+0.25, cy+1.06, 2.75, 0.65, size=12,   color=GOLD2)
    txt(sl, action,  cx+0.25, cy+2.55, 2.75, 0.38, size=11.5, bold=True, color=SAGE)

# live dot
d = box(sl, 0.9, 5.05, 0.14, 0.14, RGBColor(0x22,0xC5,0x5E), r=50000)
txt(sl, "Live", 1.15, 4.98, 1.5, 0.35, size=11.5, color=RGBColor(0x22,0xC5,0x5E))

# ══════════════════════════════════════════════════════════════════════════════
#  07 — AI TRIAGE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
box(sl, 0, 0, 13.33, 7.5, INK)

g = circ(sl, 3.5, 4.2, 4.5, RGBColor(0x5B,0x21,0xB6)); alpha(g, 8)

txt(sl, "AI Triage", 0.9, 0.55, 6, 0.45, size=11, color=GOLD)
rule(sl, 0.9, 1.1, 0.6)
txt(sl, "Claude understands\ncontext, not keywords.",
    0.9, 1.25, 8, 1.65, size=48, bold=True, color=WHITE)
txt(sl, "Every request classified. Every decision logged.",
    0.9, 3.05, 8, 0.5, size=16, color=GOLD2)

categories = [
    (RGBColor(0xEF,0x44,0x44), "Urgent Medical"),
    (RGBColor(0xF9,0x73,0x16), "In-Person Caregiver"),
    (GOLD,                     "Routine Comfort"),
    (RGBColor(0x8B,0x5C,0xF6), "Automated Task"),
    (SAGE,                     "Family Comms"),
    (MIST,                     "General Chat"),
]
for i, (col, label) in enumerate(categories):
    row, c = divmod(i, 3)
    ex = 0.9 + c * 4.1
    ey = 3.85 + row * 1.45
    cb = box(sl, ex, ey, 3.8, 1.2, INK2, r=30000); shadow(cb)
    rule(sl, ex, ey, 3.8, col, h=0.04)
    txt(sl, label, ex+0.3, ey+0.28, 3.3, 0.65, size=15, bold=True, color=WHITE)

txt(sl, '"False positives acceptable. False negatives are not."',
    0.9, 7.02, 12, 0.38, size=11.5, italic=True, color=MIST)

# ══════════════════════════════════════════════════════════════════════════════
#  08 — AUTOMATION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
box(sl, 0, 0, 13.33, 7.5, INK)

g = circ(sl, 10.5, 5.5, 4.5, SAGE); alpha(g, 7)

txt(sl, "Automation", 0.9, 0.55, 6, 0.45, size=11, color=GOLD)
rule(sl, 0.9, 1.1, 0.6)
txt(sl, "AI that acts,\nnot just advises.",
    0.9, 1.25, 7, 1.65, size=52, bold=True, color=WHITE)
txt(sl, "Simular executes desktop tasks directly — no caregiver required.",
    0.9, 3.05, 7, 0.5, size=16, color=GOLD2)

# Example commands
examples = [
    "Play Frank Sinatra",
    "Call my daughter",
    "Look up tonight's weather",
    "Find a chicken soup recipe",
    "Send a note to the front desk",
]
for i, ex in enumerate(examples):
    ey = 3.85 + i * 0.6
    b = box(sl, 0.9, ey, 5.8, 0.48, INK2, r=25000)
    txt(sl, ex, 1.2, ey+0.05, 5.3, 0.38, size=14, color=CREAM)

# Flow on right
txt(sl, "How it works", 7.5, 0.55, 5.5, 0.45, size=11, color=MIST)
steps = [
    ("Voice request arrives",       "Patient speaks their need"),
    ("Claude classifies as digital","Automated task identified"),
    ("Script generated",            "Simulang TypeScript written"),
    ("Patient confirms",            "AI speaks the task aloud"),
    ("Simular runs it",             "Task executed on desktop"),
]
for i, (title, sub) in enumerate(steps):
    sy = 1.2 + i * 1.2
    num = box(sl, 7.5, sy+0.08, 0.42, 0.42, GOLD, r=50000)
    txt(sl, str(i+1), 7.5, sy+0.1, 0.42, 0.3, size=12, bold=True,
        color=INK, align=PP_ALIGN.CENTER)
    txt(sl, title, 8.1, sy,      5.0, 0.38, size=13.5, bold=True, color=WHITE)
    txt(sl, sub,   8.1, sy+0.4,  5.0, 0.38, size=12, color=MIST)
    if i < len(steps)-1:
        rule(sl, 7.69, sy+0.55, 0.04, MIST, h=0.5)

# ══════════════════════════════════════════════════════════════════════════════
#  09 — FAMILY + MEMORY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
box(sl, 0, 0, 13.33, 7.5, INK)

txt(sl, "Family & Memory", 0.9, 0.55, 6, 0.45, size=11, color=GOLD)
rule(sl, 0.9, 1.1, 0.6)
txt(sl, "Families stay\nconnected.",
    0.9, 1.25, 7, 1.5, size=52, bold=True, color=WHITE)

# Two panels
panels = [
    (GOLD,  "Family Summaries",
     "Claude writes a daily digest — warm, human, and personal — sent to family members.",
     ["Positive moments highlighted", "Concerns flagged proactively", "Not robotic. Not clinical."]),
    (SAGE,  "Patient Memory",
     "Redis stores every interaction. Claude knows preferences and history before speaking.",
     ["Short-term session context", "Long-term preferences", "Richer triage every time"]),
]
for i, (accent, title, desc, bullets) in enumerate(panels):
    px = 0.9 + i * 6.2
    p = box(sl, px, 3.1, 5.8, 3.9, INK2, r=40000); shadow(p)
    rule(sl, px, 3.1, 5.8, accent, h=0.04)
    txt(sl, title, px+0.3, 3.3,  5.2, 0.5,  size=19, bold=True, color=WHITE)
    txt(sl, desc,  px+0.3, 3.88, 5.2, 0.95, size=13, color=GOLD2)
    for j, b in enumerate(bullets):
        rule(sl, px+0.3, 5.1+j*0.55, 0.2, accent, h=0.28)
        txt(sl, b, px+0.65, 5.05+j*0.55, 4.7, 0.42, size=12.5, color=CREAM)

# sample quote
sq = box(sl, 0.9, 7.0, 11.55, 0.38, RGBColor(0x14,0x14,0x22), r=20000)
txt(sl,
    '"Mary had a wonderful Friday — she asked for Sinatra and shared memories of her garden. A calm, happy day."',
    1.05, 7.04, 11.2, 0.3, size=10.5, italic=True, color=MIST)

# ══════════════════════════════════════════════════════════════════════════════
#  10 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
box(sl, 0, 0, 13.33, 7.5, INK)

txt(sl, "Stack", 0.9, 0.55, 6, 0.45, size=11, color=GOLD)
rule(sl, 0.9, 1.1, 0.6)
txt(sl, "Best-in-class,\nend to end.",
    0.9, 1.25, 9, 1.5, size=52, bold=True, color=WHITE)

cols = [
    ("Frontend",   ["Next.js 14", "TypeScript", "Electron", "WebSocket"]),
    ("Backend",    ["FastAPI", "Pydantic v2", "Python 3.12", "Uvicorn"]),
    ("AI & Data",  ["Claude Sonnet", "Deepgram STT", "Redis", "Arize"]),
    ("Automation", ["Simular", "Simulang", "Sentry", "Electron Builder"]),
]
for i, (cat, items) in enumerate(cols):
    cx = 0.85 + i * 3.1
    txt(sl, cat, cx, 3.1, 2.8, 0.4, size=11, bold=True, color=GOLD)
    rule(sl, cx, 3.55, 2.65, GOLD, h=0.02)
    for j, item in enumerate(items):
        iy = 3.75 + j * 0.88
        ib = box(sl, cx, iy, 2.8, 0.72, INK2, r=25000); shadow(ib)
        txt(sl, item, cx+0.2, iy+0.17, 2.4, 0.4, size=14, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
#  11 — MARKET
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
box(sl, 0, 0, 13.33, 7.5, INK)

g = circ(sl, 7.5, 4.5, 5.5, GOLD); alpha(g, 5)

txt(sl, "Market", 0.9, 0.55, 6, 0.45, size=11, color=GOLD)
rule(sl, 0.9, 1.1, 0.6)
txt(sl, "$1.4 trillion.",
    0.9, 1.25, 9, 1.3, size=72, bold=True, color=WHITE)
txt(sl, "Global elder care market — and today's tools are call buttons and clipboards.",
    0.9, 2.7, 9.5, 0.65, size=18, color=GOLD2)

stats = [
    ("$1.4T",  "Global market size"),
    ("54M+",   "Americans needing daily care"),
    ("3.7M",   "Nursing facility staff"),
    ("10,000", "Boomers turning 65 daily"),
]
for i, (val, lbl) in enumerate(stats):
    sx = 0.9 + i * 3.1
    sb = box(sl, sx, 3.65, 2.8, 2.1, INK2, r=35000); shadow(sb)
    txt(sl, val, sx+0.2, 3.78, 2.4, 0.9,  size=36, bold=True, color=GOLD)
    txt(sl, lbl, sx+0.2, 4.7,  2.4, 0.65, size=13, color=GOLD2)

txt(sl, "The problem is proven. The infrastructure is ready. The moment is now.",
    0.9, 6.15, 12, 0.5, size=14, italic=True, color=MIST, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  12 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
box(sl, 0, 0, 13.33, 7.5, INK)

g1 = circ(sl, 6.67, 3.75, 5.5, GOLD); alpha(g1, 6)
g2 = circ(sl, 6.67, 3.75, 2.5, GOLD); alpha(g2, 10)

rule(sl, 0, 0, 13.33, GOLD, h=0.06)

txt(sl, "tendly",
    0.7, 1.5, 12, 1.9, size=96, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rule(sl, 4.17, 3.55, 5.0)

txt(sl, "AI-Powered Care for Every Voice",
    0.7, 3.8, 12, 0.65, size=22, color=GOLD2, align=PP_ALIGN.CENTER)

txt(sl,
    "Built with Claude · Deepgram · Simular · Redis · Arize · Next.js · FastAPI · Electron",
    0.7, 5.4, 12, 0.45, size=12, color=MIST, align=PP_ALIGN.CENTER)

txt(sl, "Hackathon Demo  ·  2025",
    0.7, 6.7, 12, 0.4, size=11, color=DIM, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/henrybonomolo/Documents/Codework/tendly/Tendly_Pitch_Deck.pptx"
prs.save(out)
print(f"Saved → {out}  ({len(prs.slides)} slides)")
